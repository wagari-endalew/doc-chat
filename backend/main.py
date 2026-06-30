"""
DocChat RAG Backend — FastAPI + Gemini + NumPy Vector Store
Supports PDF, DOCX, TXT, PPTX | Amharic & English
No compiled C-extension dependencies (pure Python + NumPy).
"""

import os, io, uuid, json, asyncio
from pathlib import Path
from typing import List, Optional

from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import StreamingResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

import numpy as np

# Document parsers
import pypdf
import docx
from pptx import Presentation

# Embeddings
import google.generativeai as genai

# ── Config ────────────────────────────────────────────────────────────────────
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "AQ.Ab8RN6IBV_V0wy7stwGuH2wXC8BIAsWlKV0udoLyvjjlgrXTbA")
genai.configure(api_key=GEMINI_API_KEY)

CHUNK_SIZE    = 800
CHUNK_OVERLAP = 150
TOP_K         = 5

# ── Init ──────────────────────────────────────────────────────────────────────
app = FastAPI(title="DocChat API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

gemini_model = genai.GenerativeModel("gemini-flash-latest")
embed_model  = "models/gemini-embedding-001"

# ── In-memory vector store (pure Python/NumPy, no compiled deps) ───────────────
class VectorStore:
    def __init__(self):
        self.ids: List[str] = []
        self.vectors: List[np.ndarray] = []
        self.documents: List[str] = []
        self.metadatas: List[dict] = []

    def add(self, ids, embeddings, documents, metadatas):
        for i, e, d, m in zip(ids, embeddings, documents, metadatas):
            self.ids.append(i)
            self.vectors.append(np.array(e, dtype=np.float32))
            self.documents.append(d)
            self.metadatas.append(m)

    def count(self) -> int:
        return len(self.ids)

    def query(self, query_embedding: List[float], n_results: int = 5):
        if not self.vectors:
            return {"documents": [[]], "metadatas": [[]]}

        q = np.array(query_embedding, dtype=np.float32)
        q_norm = q / (np.linalg.norm(q) + 1e-10)

        mat = np.stack(self.vectors)
        mat_norm = mat / (np.linalg.norm(mat, axis=1, keepdims=True) + 1e-10)

        sims = mat_norm @ q_norm
        top_idx = np.argsort(-sims)[:n_results]

        docs  = [self.documents[i] for i in top_idx]
        metas = [self.metadatas[i] for i in top_idx]
        return {"documents": [docs], "metadatas": [metas]}

    def clear(self):
        self.ids, self.vectors, self.documents, self.metadatas = [], [], [], []

    def get_all_metadatas(self):
        return self.metadatas


store = VectorStore()

# ── Document extraction ─────────────────────────────────────────────────────────
def extract_text_pdf(data: bytes) -> str:
    reader = pypdf.PdfReader(io.BytesIO(data))
    return "\n".join(p.extract_text() or "" for p in reader.pages)

def extract_text_docx(data: bytes) -> str:
    doc = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs)

def extract_text_txt(data: bytes) -> str:
    return data.decode("utf-8", errors="ignore")

def extract_text_pptx(data: bytes) -> str:
    prs  = Presentation(io.BytesIO(data))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

EXTRACTORS = {
    ".pdf":  extract_text_pdf,
    ".docx": extract_text_docx,
    ".txt":  extract_text_txt,
    ".pptx": extract_text_pptx,
}

def chunk_text(text: str, size=CHUNK_SIZE, overlap=CHUNK_OVERLAP) -> List[str]:
    chunks, start = [], 0
    while start < len(text):
        end = start + size
        chunks.append(text[start:end])
        start += size - overlap
    return [c.strip() for c in chunks if c.strip()]

def embed_texts(texts: List[str]) -> List[List[float]]:
    result = genai.embed_content(
        model=embed_model,
        content=texts,
        task_type="retrieval_document",
    )
    return result["embedding"]

def embed_query(query: str) -> List[float]:
    result = genai.embed_content(
        model=embed_model,
        content=query,
        task_type="retrieval_query",
    )
    return result["embedding"]

# ── Routes ────────────────────────────────────────────────────────────────────
@app.get("/health")
def health():
    return {"status": "ok", "docs_indexed": store.count()}

@app.post("/upload")
async def upload_documents(files: List[UploadFile] = File(...)):
    results = []
    for file in files:
        suffix = Path(file.filename).suffix.lower()
        if suffix not in EXTRACTORS:
            results.append({"file": file.filename, "status": "unsupported"})
            continue

        data = await file.read()
        try:
            text   = EXTRACTORS[suffix](data)
            chunks = chunk_text(text)
            if not chunks:
                results.append({"file": file.filename, "status": "empty"})
                continue

            ids, embeddings, metas, docs = [], [], [], []
            for i in range(0, len(chunks), 100):
                batch = chunks[i:i+100]
                embs  = embed_texts(batch)
                for j, (chunk, emb) in enumerate(zip(batch, embs)):
                    ids.append(str(uuid.uuid4()))
                    embeddings.append(emb)
                    docs.append(chunk)
                    metas.append({"source": file.filename, "chunk": i + j})

            store.add(ids=ids, embeddings=embeddings, documents=docs, metadatas=metas)
            results.append({"file": file.filename, "status": "ok", "chunks": len(chunks)})
        except Exception as e:
            results.append({"file": file.filename, "status": "error", "detail": str(e)})

    return {"results": results}

class QueryRequest(BaseModel):
    question: str
    language: str = "auto"

@app.post("/chat/stream")
async def chat_stream(req: QueryRequest):
    q_emb = embed_query(req.question)
    hits  = store.query(q_emb, n_results=min(TOP_K, store.count() or 1))

    context_chunks = hits["documents"][0] if hits["documents"] else []
    sources        = list({m["source"] for m in hits["metadatas"][0]}) if hits["metadatas"][0] else []
    context        = "\n\n---\n\n".join(context_chunks)

    lang_hint = {
        "am":   "ምላሽህን በአማርኛ ስጥ።",
        "en":   "Respond in English.",
        "auto": "Detect the language of the question and respond in the same language (Amharic or English).",
    }.get(req.language, "auto")

    system_prompt = f"""You are an intelligent document assistant that supports both Amharic (አማርኛ) and English.
{lang_hint}

Use ONLY the context below to answer. If the answer is not in the context, say so clearly.
Cite the source document name when relevant.

CONTEXT:
{context}
"""

    async def event_generator():
        try:
            response = gemini_model.generate_content(
                [system_prompt, req.question],
                stream=True,
                generation_config=genai.types.GenerationConfig(
                    temperature=0.3,
                    max_output_tokens=2048,
                ),
            )
            for chunk in response:
                if chunk.text:
                    payload = json.dumps({"token": chunk.text, "sources": sources})
                    yield f"data: {payload}\n\n"
                    await asyncio.sleep(0)

            yield f"data: {json.dumps({'done': True, 'sources': sources})}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )

@app.delete("/documents")
def clear_documents():
    store.clear()
    return {"status": "cleared"}

@app.get("/documents")
def list_documents():
    if store.count() == 0:
        return {"documents": []}
    sources = list({m["source"] for m in store.get_all_metadatas()})
    return {"documents": sources, "total_chunks": store.count()}